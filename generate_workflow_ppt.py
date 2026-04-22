import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Color Palette
DARK_NAVY = RGBColor(44, 62, 80)
BLUE_HEADER = RGBColor(41, 128, 185)
TEXT_MAIN = RGBColor(44, 62, 80)
TEXT_SUB = RGBColor(60, 60, 60)
TEXT_LIGHT = RGBColor(100, 100, 100)
ACCENT_RED = RGBColor(192, 57, 43)

def apply_slide_theme(slide, title_text, is_title_slide=False):
    shapes = slide.shapes
    if is_title_slide:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_NAVY
        
        title_shape = shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(2))
        tf = title_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = "微软雅黑"
        p.font.size = Pt(48)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        line = shapes.add_shape(1, Inches(4), Inches(3.6), Inches(2), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(52, 152, 219)
        line.line.fill.background()
    else:
        header_rect = shapes.add_shape(
            1, Inches(0), Inches(0), Inches(10), Inches(1.1)
        )
        header_rect.fill.solid()
        header_rect.fill.fore_color.rgb = BLUE_HEADER
        header_rect.line.fill.background()
        
        title_shape = shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.8))
        tf = title_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = "微软雅黑"
        p.font.size = Pt(32)
        p.font.bold = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_content(slide, bullets_dict, font_size_offset=0):
    shapes = slide.shapes
    body_shape = shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.8))
    tf = body_shape.text_frame
    tf.word_wrap = True
    
    first = True
    for level, text in bullets_dict:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            
        p.text = text
        p.level = level
        p.font.name = "微软雅黑"
        p.line_spacing = 1.15  
        
        if level == 0:
            p.font.size = Pt(22 + font_size_offset)
            p.font.bold = True
            p.font.color.rgb = TEXT_MAIN
            p.space_before = Pt(14)
            p.space_after = Pt(4)
        elif level == 1:
            p.font.size = Pt(16 + font_size_offset)
            p.font.color.rgb = TEXT_SUB
            p.space_before = Pt(6)
            p.space_after = Pt(2)
        else:
            p.font.size = Pt(14 + font_size_offset)
            p.font.color.rgb = TEXT_LIGHT
            p.space_before = Pt(4)

def add_comparison_table(slide):
    shapes = slide.shapes
    rows = 4
    cols = 4
    left = Inches(0.5)
    top = Inches(2.2)
    width = Inches(9.0)
    height = Inches(3.0)
    
    table_shape = shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(2.5)
    
    headers = ["对比维度", "真实编译", "Tree-sitter\n(静态语法解析)", "LSP 工具链\n(OS-Agent D)"]
    for i, title in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = title
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(41, 128, 185)
        p.alignment = PP_ALIGN.CENTER
        
    data = [
        ["编译流水线深度", "预处理 -> 语法 -> 语义 -> IR代码 -> 寄存器分配/汇编 -> 链接", "止步于 语法分析\n(无法处理宏展开)", "深入至 语义分析\n(支持跨文件符号消解与类型推导)"],
        ["函数调用准确率", "100%", "容易产生误判\n(纯文本正则匹配，无法区分同名函数、宏与指针)", "高度精确\n(基于语义分析处理宏替换与结构定义)"],
        ["工程实施优劣", "依赖交叉编译环境与 Makefile，在 OS 库中极易因依赖缺失而中断。", "解析速度快且免配置，适合宏观特征提取。但分析深度不足。", "免除汇编环节的硬环境要求，保留语义级的代码追溯能力，符合静态分析需求。"]
    ]
    
    for row_idx, row_data in enumerate(data):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.font.name = "微软雅黑"
            p.font.size = Pt(12)
            if col_idx == 0:
                p.font.bold = True
                p.font.color.rgb = TEXT_MAIN
            else:
                p.font.color.rgb = TEXT_SUB
            p.alignment = PP_ALIGN.CENTER if col_idx == 0 else PP_ALIGN.LEFT

def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # === 1. Title ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(slide, "OS-Agent-D 工作流程剖析", is_title_slide=True)
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(8), Inches(1.5))
    tf = subtitle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "解析 os_agent_d_describe.py 的自动化代码分析管线设计"
    p.font.color.rgb = RGBColor(200, 214, 229)
    p.font.name = "微软雅黑"
    p.font.size = Pt(22)
    p.alignment = PP_ALIGN.CENTER

    # === 2. 整体管线宏观视图 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(slide, "第一部分：管线宏观视图与结构导读")
    add_content(slide, [
        (0, "系统管线划分为四大阶段："),
        (1, "➤ 一、基础设施构建 (阶段0/0.5)：包含 Git 代码同步、Repo Profile 特征画像、RAG 语义块向量化及全局路径记忆池。"),
        (1, "➤ 二、工具链体系：基于 LSP (Language Server Protocol) 语义解析的底层探针系统。"),
        (1, "➤ 三、Plan-Execute-Review 分析闭环：拆分 10 大核心专题，依据三重核心 Prompt 执行严格的推理与代码取证。"),
        (1, "➤ 四、产物合成：Markdown 报告合成、首章表格矩阵生成及 Call Graph 逻辑图提取。"),
    ])

    # === 3. 基础设施 1：阶段0 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(slide, "基础设施：阶段 0 环境沙盒化")
    add_content(slide, [
        (0, "隔离运行环境："),
        (1, "➤ 采用 `git clone` 将目标操作系统仓库提取至 `./repos/` 沙盒目录，保证多实例并发分析无干扰。"),
        (0, "分析状态重置："),
        (1, "➤ 强制清理分析遗留数据，包括 Rust 的 target 目录、LSP 编译缓存数据库等，消除跨项目的环境污染。"),
    ])

    # === 4. 基础设施 2：Repo Profile ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(slide, "基础设施：Repo Profile 静态特征扫描")
    add_content(slide, [
        (0, "静态启发式分析："),
        (1, "➤ 在 LLM 介入前，系统使用 Python 脚本执行全局静态扫描 (`build_repo_profile`)。"),
        (0, "具体扫描指标："),
        (1, "➤ 构建系统探查：探测根目录中的 `Makefile`, `Cargo.toml`, `CMakeLists.txt`。"),
        (1, "➤ 语言占比统计：计算 `.c`, `.h`, `.rs`, `.S` 文件分布，判定主实现语言。"),
        (1, "➤ 硬件架构推断：扫描特定目录 (如 `arch/riscv`, `arch/arm`) 判断系统指令集架构。"),
        (0, "实现作用："),
        (1, "➤ 输出 `repo_profile.json`，为后续 LLM 分析提供底层技术栈的硬先验上下文。"),
    ], font_size_offset=-1)

    # === 5. 基础设施 3：阶段 0.5 RAG ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(slide, "基础设施：阶段 0.5 RAG 向量化索引")
    add_content(slide, [
        (0, "语义检索架构："),
        (1, "➤ ① 语法感知切分 (Syntax-Aware Chunking)：不依赖代码行数截断，而是通过语法树边界 (如函数、结构体定义) 划分源码块，维持逻辑完整度。"),
        (1, "➤ ② 向量化计算 (Embedding)：利用专用的代码表征模型（如 `jina-embeddings-v2`）将切分块映射至向量空间。"),
        (1, "➤ ③ 向量存储与检索：在本地建立 Vector DB，供 Execute 阶段基于特征语义映射执行高召回率检索。"),
    ], font_size_offset=-1)

    # === 6. 基础设施 4：全局路径记忆 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(slide, "基础设施：跨阶段全局路径记忆池")
    add_content(slide, [
        (0, "跨模块耦合挑战："),
        (1, "➤ 操作系统各子系统存在底层代码耦合（如调度队列依赖内存结构）。若各阶段隔离，极易引发冗余代码搜索。"),
        (0, "缓存池实现："),
        (1, "➤ 数据抽取：子阶段结束后，正则表达式提取已论证的核心源码路径（如 `kernel/mm.c:45`）。"),
        (1, "➤ 队列维护：合并至全局路径池，并执行去重操作，上限保持 100 条热点代码索引。"),
        (1, "➤ 上下文传递：在后续阶段的 Plan 环节，将路径池作为先验信息注入 Prompt，收敛 LLM 的搜索范围。"),
    ], font_size_offset=-1)

    # === 7. 工具链对比表 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(slide, "第二部分：代码静态分析工具横评")
    add_comparison_table(slide)

    # === 8. 全套 Tools 工具集解析 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(slide, "核心探针系统：Tools 工具链与底层依赖机制")
    add_content(slide, [
        (0, "LSP 引擎的具体实现与核心依赖："),
        (1, "➤ 引擎绑定：C/C++ 项目挂载 `clangd` 进程；Rust 项目挂载 `rust-analyzer`。"),
        (1, "➤ 核心基石：LSP 之所以能提供跨文件宏解析，必须依赖全局编译信息。系统需通过 `bear make` 等工具提取 `compile_commands.json`（或兜底构造 `compile_flags.txt`），为引擎注入包含路径 (-I) 与宏定义 (-D)。这是 LSP 能否正确消解符号的生死关键。"),
        (0, "LSP 具体下发接口："),
        (1, "➤ `lsp_get_definition` (符号定义解析)、`lsp_get_call_graph` (调用栈推理)、`lsp_get_references` (全局引用检索)。"),
        (0, "兜底检索工具组："),
        (1, "➤ `rag_search_code` (基于向量语义召回)、`grep_in_repo` (纯文本正则兜底扫描)。"),
        (0, "文件 IO 工具组："),
        (1, "➤ `list_repo_structure` (目录遍历)、`read_code_segment` (源码切片读取提取证据)。"),
    ], font_size_offset=-2)

    # === 9. P-E-R：Prompt 体系 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(slide, "第三部分：P-E-R 核心机制与 Prompt 体系")
    add_content(slide, [
        (0, "系统通过三层 Prompt 结构对大模型进行刚性约束："),
        (1, "➤ [1] System Prompt (全局底座原则)：贯穿全流程，确立大模型作为“高级逆向工程专家”的身份。禁止臆测，要求所有的断言必须有工具输出的代码行号作为绝对支撑。"),
        (1, "➤ [2] Stage Prompt (章级约束与上下文)：在每个 P-E-R 模块运行时动态注入。它包含了该章节特定的聚焦任务要求、`repo_profile.json` 画像特征，以及从之前阶段继承的“全局路径记忆池”，用于裁剪当前的搜索空间。"),
        (1, "➤ [3] JSON-QA 题单配置：针对 02 至 09 章节的专门约束，强制模型将其探索到的代码证据与结论，严格填充进预定义的 JSON Schema 中，作为后续机器自动化评阅的基石。"),
    ], font_size_offset=-1)

    # === 10. P-E-R：① Plan 阶段 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(slide, "分析管线：① Plan 阶段 (策略规划)")
    add_content(slide, [
        (0, "核心目的：探索方向，收敛范围，防止后续在海量代码中陷入盲目循环。"),
        (0, "传递给 LLM 的输入 (What goes in)："),
        (1, "➤ 全局 System Prompt + 对应章节的 Stage Prompt (包含了启发式画像库 Repo Profile 与热点路径记忆池)。"),
        (0, "模型在此阶段干了什么 (What it does)："),
        (1, "➤ 代理 (Planner) 在沙盒内自主调用底层工具执行轻量级检索探路。"),
        (1, "➤ 它负责在茫茫代码中寻找真正的功能入口点，排除伪命题与干扰项，但不编写最终正文。"),
        (0, "阶段性产出 (What comes out)："),
        (1, "➤ 形成一份锁死的 `execution_steps` (执行步骤清单契约)。要求后续环节必须按此确定的路径执行，彻底杜绝发散。"),
    ], font_size_offset=-1)

    # === 11. P-E-R：② Execute 阶段 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(slide, "分析管线：② Execute 阶段 (深度取证与编写)")
    add_content(slide, [
        (0, "核心目的：顺藤摸瓜，拿到铁证代码，撰写最终的技术分析正文。"),
        (0, "传递给 LLM 的输入 (What goes in)："),
        (1, "➤ 全局 System Prompt + Plan 阶段输出的 `execution_steps` 强制计划 + 本阶段特定的题单任务合并 (Human Message)。"),
        (0, "模型在此阶段干了什么 (What it does)："),
        (1, "➤ 搭载契约进入 ReAct 闭环循环。按照 RAG -> LSP -> Grep 的优先级顺序，严格调动全套探针深层下钻。"),
        (1, "➤ 坚守“无源码无结论”的铁律，自带 ErrorTracker 处理 LSP 编译崩溃后的退避与自愈。"),
        (0, "阶段性产出 (What comes out)："),
        (1, "➤ 输出极富证据支撑的该章节 Markdown 总结文本。"),
        (1, "➤ (对于 02-09 章) 输出按强制格式填充好的 JSON-QA 原始数据体。"),
    ], font_size_offset=-1)

    # === 12. P-E-R：③ Review 阶段 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(slide, "分析管线：③ Review 阶段 (可选的强校验)")
    add_content(slide, [
        (0, "核心目的：修补破损格式，将自然语言断言转化为三态机器逻辑。"),
        (0, "并非所有章节都参与 Review："),
        (1, "➤ 第 01 章 (概览) 与 第 10 章 (历史) 仅生成综合性 Markdown 报告，不产出具体的题单判定，故跳过此阶段。"),
        (0, "传递给 LLM 的输入 (What goes in)："),
        (1, "➤ 仅发送原始的 `describe_stage_qa` 题单模板，以及 Execute 阶段输出的未处理的 JSON 字符串（不再带入杂乱的工具摘录上下文，保持极度专注）。"),
        (0, "模型与系统在此阶段干了什么与产出："),
        (1, "➤ LLM 修正破损或非法结构的 JSON 闭合，将文本强行拉回到 `tri_state_impl` (implemented/stub/not_found) 的三态评估口径上。"),
        (1, "➤ 随后 Python 控制流 (`coerce_answers_payload`) 接管验证，强制覆写遗漏的缺失字段。确保入库数据类型绝对安全。"),
    ], font_size_offset=-2)

    # === 13. 题单体系 (全页整合) ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(slide, "结构化指标评价：全模块 JSON 题单配置库")
    add_content(slide, [
        (0, "02-09 深度专题探测 (配置 JSON 题单，对齐教材核心维度)："),
        (1, "➤ [02 启动] Trapframe 上下文数据结构、特权级切换控制流、系统调用入口。"),
        (1, "➤ [03 内存] 物理帧分配算法底座 (Buddy)、多级页表机制关联 (Sv39)、惰性分配。"),
        (1, "➤ [04 进程] PCB 数据模型、进程调度核心算法逻辑、Context Switch 保存策略。"),
        (1, "➤ [05 文件] VFS 抽象层系统调用对接、inode 生命周期、底层块设备驱动映射。"),
        (1, "➤ [06 同步] 自旋/互斥锁及关中断策略、核心进程间通信 (IPC) 结构。"),
        (1, "➤ [07 安全] User-Mode 执行边界保护、非法内存越界拦截响应逻辑。"),
        (1, "➤ [08 网络] 网卡收发包外部中断处理、协议栈分层设计结构 (TCP/UDP)。"),
        (1, "➤ [09 调试] Kernel Panic 全局流追踪、内存 Dump 落盘格式、Assert 断言机制。"),
        (0, "01与10 宏观统筹阶段 (无题单设计)："),
        (1, "➤ [01 概览] 与 [10 历史] 的职责是对全局架构特征进行高维浓缩归纳，无法也不应拆解为死板的三态技术指标，故完全移除静态题单，交由宏观 Prompt 驱动全文。"),
    ], font_size_offset=-3)

    # === 14. 第 01 章强制格式 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(slide, "第四部分：Overview 报告格式规范约束")
    add_content(slide, [
        (0, "格式契约化设计："),
        (1, "➤ 第一章负责输出全盘系统评估状态。为满足后续评测程序的自动解析需求，系统在 Prompt 层面部署硬性排版规范："),
        (1, "➤ 静态表格矩阵约束：要求强制输出 10 行定长宽的《子系统完成度矩阵》。禁止 LLM 依据各子系统实现规模而删减或扩写表格行。"),
        (1, "➤ 降级标记限制：规定必须且只能使用特定的 4 级 Markdown 标题（“技术清单”、“关键实现”、“依赖”、“模块衔接”）。"),
        (1, "➤ 字符串锚定：要求技术栈说明段落前缀文字需 100% 对齐预设模板（常量对齐）。"),
        (0, "工程价值："),
        (1, "➤ 规避 LLM 生成长文本时带来的“排版随机性”，保障任意代码库产生的报告结构具备严格对齐特性。"),
    ], font_size_offset=-2)

    # === 15. Call Graph 生成 (全面重构回归工具链源码实现) ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(slide, "全局架构可视化：Domain × Layer 拓扑大图生成")
    add_content(slide, [
        (0, "底层算法执行链 (`callgraph_overview.py`)："),
        (1, "➤ [1] 初筛建树：利用 Tree-sitter 基于语法树提取全仓所有的函数与调用流，生成基础有向图 (DiGraph)。"),
        (1, "➤ [2] PageRank 过滤：提取海量函数后，系统计算入度与出度，利用 PageRank 权重算法排出全局最重要的 Top-K 核心调度枢纽，剔除海量底层噪音。"),
        (1, "➤ [3] LSP 边细化：对过滤出的核心节点，唤醒 LSP 的 `callHierarchy` 重新修正跨文件依赖，消灭因同名宏与函数指针导致的误判伪连线。"),
        (0, "领域划分与动态 SVG 渲染："),
        (1, "➤ 二维矩阵聚类：调用 LLM 对核心函数进行批处理分类，纵轴为运行层级 (Userspace/Syscall/Kernel/Hardware)，横轴为功能域 (Trap/Mem/Sched 等)。"),
        (1, "➤ 动态 SVG 连线排版：提取连线逻辑，脚本自动根据行列节点密度动态计算、自适应调整表格的宽高比例，生成一张跨域节点连线清晰、具备色彩辨识度的高清 SVG 全景图，嵌入在评测报告首页。"),
    ], font_size_offset=-3)

    # === 16. 设想：OS-Agent C 模块与机制 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(slide, "架构前瞻：OS-Agent C 评测横评设计 (1)")
    add_content(slide, [
        (0, "基于 Agent D 数据的结构化对齐评估："),
        (0, "设计评测机制示例 (待落地的抽象评估模块)："),
        (1, "➤ [内存管理模块] 核心比对：物理内存底座结构（如 OS-A Buddy vs OS-B Bitmap）；多级页表硬件寄存器配置差异（Sv39 vs Sv48）。"),
        (1, "➤ [进程调度模块] 核心比对：进程抽象模型（独立进程 vs 多线程共享机制）；时钟调度算法复杂度与实现方式（O(1) 优先级 vs CFS 公平调度）。"),
        (1, "➤ [文件系统模块] 核心比对：VFS 系统调用接口标准化程度；底层的 Inode 缓存映射逻辑差异。"),
    ], font_size_offset=-1)

    # === 17. OS Agent C 与 D 题单反哺 (人工闭环修订) ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(slide, "架构前瞻：OS-Agent C 粗筛维度与人工反馈机制 (2)")
    add_content(slide, [
        (0, "免 LLM 的特征粗筛 (Coarse Filtering)："),
        (1, "➤ 在架构设计上，Agent C 首先利用轻量化脚本提取各模块的关键机制维度 (Key)，通过哈希比对两个 OS 的 JSON 产物。这保证了极高吞吐量与零算力消耗。"),
        (0, "面临的同构框架挑战："),
        (1, "➤ 在早期的初步测试阶段，大量的学生级 OS 源于同一教学框架变体（如 rCore, ArceOS, uCore 等）。"),
        (1, "➤ 若粗筛提取的维度不够细致，将导致“区分度过低”，无法判定这究竟是合理的架构修改还是高度同质化。"),
        (0, "人工测试反馈与逆向重写闭环："),
        (1, "➤ 鉴于上述挑战，在 Agent C 粗筛测试运行过程中，测试人员主动通过观察评测结果来发现系统层面的鉴别盲区。"),
        (1, "➤ 人工介入提炼更细粒度的鉴别特征，将这些深层次技术点（例如锁粒度差异、特殊任务队列）转化为新的 Key，并注入回 Agent D 的 `stage_qa.json` 题单。"),
        (1, "➤ 由此形成了“实测同构框架的区分度瓶颈，驱动人工精细化迭代底层题库”的技术进化闭环。"),
    ], font_size_offset=-2)

    # === 18. 探讨与 Q&A ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(slide, "架构定调探讨")
    
    shapes = slide.shapes
    body_shape = shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(3))
    tf = body_shape.text_frame
    tf.word_wrap = True
    
    p2 = tf.paragraphs[0]
    p2.text = "关于 OS-Agent C 架构的评估标的：\n\n后续研发投入的最高优先级，\n应侧重于操作系统的“架构横向评比设计 (Comparison)”？\n还是倾向于微观源码维度的“抄袭判定鉴查 (Plagiarism Detection)”？"
    p2.font.name = "微软雅黑"
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_RED
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)
    p2.line_spacing = 1.3

    prs.save("os_agent_d_workflow_v21.pptx")
    print("PPT 已成功重构排版：os_agent_d_workflow_v21.pptx")

if __name__ == "__main__":
    create_ppt()
